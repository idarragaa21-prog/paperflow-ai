# -*- coding: utf-8 -*-
"""Helpers de layout para a apresentação (16:9, fundo branco, estilo clínico)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# Paleta
INK     = RGBColor(0x12, 0x32, 0x4F)
PRIMARY = RGBColor(0x1F, 0x6F, 0xB2)
TEAL    = RGBColor(0x0E, 0x9A, 0xA7)
ALERT   = RGBColor(0xD6, 0x45, 0x41)
AMBER   = RGBColor(0xE1, 0xA1, 0x40)
GREEN   = RGBColor(0x2E, 0x9E, 0x6B)
TXT     = RGBColor(0x45, 0x5A, 0x6B)
MUTE    = RGBColor(0x8A, 0x9B, 0xA8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF3, 0xF6, 0xF9)
LINE    = RGBColor(0xDD, 0xE4, 0xEA)
CARDBG  = RGBColor(0xF7, 0xFA, 0xFC)

FONT = "Calibri"
FONT_L = "Calibri Light"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of dicts {t, size, color, bold, font, space_after, italic, space_before}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if "space_after" in ln:
            p.space_after = Pt(ln["space_after"])
        if "space_before" in ln:
            p.space_before = Pt(ln["space_before"])
        if "line_spacing" in ln:
            p.line_spacing = ln["line_spacing"]
        runs = ln["t"] if isinstance(ln["t"], list) else [ln]
        for rr in runs:
            r = p.add_run()
            r.text = rr["t"] if isinstance(rr, dict) else ln["t"]
            f = r.font
            f.size = Pt(rr.get("size", ln.get("size", 16)) if isinstance(rr, dict) else ln.get("size", 16))
            f.bold = rr.get("bold", ln.get("bold", False)) if isinstance(rr, dict) else ln.get("bold", False)
            f.italic = rr.get("italic", ln.get("italic", False)) if isinstance(rr, dict) else ln.get("italic", False)
            f.name = rr.get("font", ln.get("font", FONT)) if isinstance(rr, dict) else ln.get("font", FONT)
            col = rr.get("color", ln.get("color", TXT)) if isinstance(rr, dict) else ln.get("color", TXT)
            f.color.rgb = col
    return tb


def header(slide, kicker, title, idx, total, accent=PRIMARY):
    """Cabeçalho padrão de slide de conteúdo."""
    # barra de acento vertical
    rect(slide, Inches(0.6), Inches(0.62), Inches(0.12), Inches(0.9), accent)
    txt(slide, Inches(0.86), Inches(0.55), Inches(11.6), Inches(0.4),
        [{"t": kicker.upper(), "size": 12.5, "color": accent, "bold": True}])
    txt(slide, Inches(0.86), Inches(0.86), Inches(11.9), Inches(0.75),
        [{"t": title, "size": 27, "color": INK, "bold": True, "font": FONT}])
    footer(slide, idx, total)


def footer(slide, idx, total, accent=PRIMARY):
    rect(slide, Inches(0.6), Inches(6.95), Inches(12.13), Pt(0.9), LINE)
    txt(slide, Inches(0.6), Inches(7.02), Inches(9.0), Inches(0.35),
        [{"t": "Artrite Séptica — do básico ao avançado", "size": 9.5, "color": MUTE}])
    txt(slide, Inches(10.5), Inches(7.02), Inches(2.23), Inches(0.35),
        [{"t": f"{idx:02d} / {total:02d}", "size": 9.5, "color": MUTE, "bold": True,
          "align": PP_ALIGN.RIGHT}])


def citation(slide, text, y=6.55):
    txt(slide, Inches(0.86), Inches(y), Inches(11.6), Inches(0.35),
        [{"t": text, "size": 9.5, "color": MUTE, "italic": True}])


def bullet_card(slide, x, y, w, h, icon_color, head, body, num=None):
    """Cartão com título e corpo curto."""
    card = rect(slide, x, y, w, h, CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.06)
    rect(slide, x, y, Inches(0.09), h, icon_color)
    tb = slide.shapes.add_textbox(x + Inches(0.28), y + Inches(0.16), w - Inches(0.45), h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_right = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = head
    r.font.size = Pt(14.5); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = INK
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(11.5); r2.font.name = FONT; r2.font.color.rgb = TXT
    return card


def _set_corner(shape, frac):
    try:
        shape.adjustments[0] = frac
    except Exception:
        pass


def stat(slide, x, y, w, big, label, color=PRIMARY, sub=None):
    lines = [{"t": big, "size": 40, "color": color, "bold": True, "align": PP_ALIGN.CENTER,
              "space_after": 2},
             {"t": label, "size": 12.5, "color": TXT, "align": PP_ALIGN.CENTER}]
    if sub:
        lines.append({"t": sub, "size": 10, "color": MUTE, "align": PP_ALIGN.CENTER, "space_before": 2})
    txt(slide, x, y, w, Inches(1.6), lines, anchor=MSO_ANCHOR.TOP)


def node(slide, x, y, w, h, text, fill, txt_color=WHITE, size=12.5, bold=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None):
    sp = rect(slide, x, y, w, h, fill, shape=shape, line=line, line_w=1.25)
    _set_corner(sp, 0.12)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(5); tf.margin_right = Pt(5); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for j, part in enumerate(text.split("\n")):
        pp = p if j == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = part
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = txt_color
    return sp


def connect(slide, x1, y1, x2, y2, color=MUTE, width=1.6, arrow=True):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if arrow:
        lnEl = cn.line._get_or_add_ln()
        tail = lnEl.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        lnEl.append(tail)
    return cn


def edge_label(slide, x, y, w, text, color=TXT):
    txt(slide, x, y, w, Inches(0.3),
        [{"t": text, "size": 10, "color": color, "bold": True, "align": PP_ALIGN.CENTER}])
