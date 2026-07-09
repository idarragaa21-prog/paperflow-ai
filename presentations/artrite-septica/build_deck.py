# -*- coding: utf-8 -*-
"""
Artrite Séptica — do básico ao avançado
Apresentação profissional (pt-BR), 16:9, fundo branco.
Autor: Diego Alejandro Idárraga — R2 Ortopedia e Traumatologia
Hospital Municipal Barata Ribeiro — Rio de Janeiro
"""
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from deck_helpers import (
    new_deck, blank, rect, txt, header, footer, citation, bullet_card, stat,
    node, connect, edge_label, _set_corner,
    INK, PRIMARY, TEAL, ALERT, AMBER, GREEN, TXT, MUTE, WHITE, LIGHT, LINE,
    CARDBG, FONT, FONT_L,
)

A = os.path.join(os.path.dirname(__file__), "assets")
prs = new_deck()
TOTAL = 30
I = Inches


def RGBColor_soft(c, amt=0.90):
    """Tinte claro de uma cor (mistura com branco) para fundos de cartão."""
    return RGBColor(int(c[0] + (255 - c[0]) * amt),
                    int(c[1] + (255 - c[1]) * amt),
                    int(c[2] + (255 - c[2]) * amt))


def img(slide, name, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = I(w)
    if h: kw["height"] = I(h)
    return slide.shapes.add_picture(os.path.join(A, name), I(x), I(y), **kw)


def table(slide, x, y, w, rows, col_w, header_fill=INK, font=10.5, row_h=0.34, head_h=0.4):
    nrows = len(rows); ncols = len(rows[0])
    gt = slide.shapes.add_table(nrows, ncols, I(x), I(y), I(w), I(head_h + row_h*(nrows-1))).table
    gt.first_row = False; gt.horz_banding = False
    for j, cw in enumerate(col_w):
        gt.columns[j].width = I(cw)
    gt.rows[0].height = I(head_h)
    for i in range(1, nrows):
        gt.rows[i].height = I(row_h)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Pt(6); c.margin_right = Pt(6)
            c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else CARDBG
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(font if i else font+0.5)
            r.font.name = FONT
            r.font.bold = (i == 0)
            r.font.color.rgb = WHITE if i == 0 else INK if j == 0 else TXT
    return gt


# ======================================================================
# SLIDE 1 — PORTADA
# ======================================================================
s = blank(prs)
rect(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)
# faixa lateral de acento
rect(s, 0, 0, I(0.35), prs.slide_height, INK)
rect(s, I(0.35), 0, I(0.08), prs.slide_height, TEAL)
# etiqueta superior
txt(s, I(1.0), I(1.15), I(11), I(0.5),
    [{"t": "REVISÃO CLÍNICA  ·  ORTOPEDIA E TRAUMATOLOGIA", "size": 13, "color": TEAL, "bold": True}])
# título
txt(s, I(0.98), I(1.75), I(11.6), I(2.2),
    [{"t": "Artrite Séptica", "size": 60, "color": INK, "bold": True, "font": FONT, "space_after": 4},
     {"t": "Do básico ao avançado — diagnóstico e manejo baseados em evidência",
      "size": 21, "color": PRIMARY, "font": FONT_L}])
# regra
rect(s, I(1.0), I(4.35), I(4.2), Pt(2.2), TEAL)
# autor / instituição
txt(s, I(1.0), I(4.65), I(11), I(1.6),
    [{"t": "Diego Alejandro Idárraga", "size": 20, "color": INK, "bold": True, "space_after": 3},
     {"t": "R2 — Ortopedia e Traumatologia", "size": 14.5, "color": TXT, "space_after": 1},
     {"t": "Hospital Municipal Barata Ribeiro — Rio de Janeiro", "size": 14.5, "color": TXT}])
# ícone-alvo (articulação estilizada) no canto direito
cx, cy = I(10.7), I(2.55)
rect(s, cx, cy, I(1.7), I(1.7), LIGHT, shape=MSO_SHAPE.OVAL)
o2 = s.shapes.add_shape(MSO_SHAPE.OVAL, cx+I(0.32), cy+I(0.32), I(1.06), I(1.06))
o2.fill.solid(); o2.fill.fore_color.rgb = WHITE; o2.line.color.rgb = ALERT; o2.line.width = Pt(2.5); o2.shadow.inherit = False
o3 = s.shapes.add_shape(MSO_SHAPE.OVAL, cx+I(0.62), cy+I(0.62), I(0.46), I(0.46))
o3.fill.solid(); o3.fill.fore_color.rgb = ALERT; o3.line.fill.background(); o3.shadow.inherit = False
txt(s, I(1.0), I(6.7), I(11), I(0.4),
    [{"t": "Rio de Janeiro · 2026", "size": 12, "color": MUTE}])


# ======================================================================
# SLIDE 2 — AGENDA
# ======================================================================
s = blank(prs)
header(s, "Roteiro", "Agenda", 2, TOTAL, accent=TEAL)
items = [
    ("01", "Fundamentos", "Definição, epidemiologia e classificação"),
    ("02", "Microbiologia e fisiopatologia", "Agentes etiológicos e vias de infecção"),
    ("03", "Quadro clínico", "Apresentação, fatores de risco e articulações"),
    ("04", "Diagnóstico", "Laboratório, líquido sinovial, imagem e algoritmos"),
    ("05", "Artrite séptica pediátrica", "Critérios de Kocher e Caird"),
    ("06", "Tratamento", "Antibioticoterapia e drenagem articular"),
    ("07", "Cenários avançados", "Artrite gonocócica e infecção de prótese"),
    ("08", "Complicações e prognóstico", "Desfechos e mensagens-chave"),
]
x0, y0 = 0.86, 1.85
cw, ch, gx, gy = 5.75, 1.05, 0.35, 0.16
for k, (num, h, b) in enumerate(items):
    col = k % 2; rowi = k // 2
    x = x0 + col*(cw+gx); y = y0 + rowi*(ch+gy)
    card = rect(s, I(x), I(y), I(cw), I(ch), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.09)
    txt(s, I(x+0.22), I(y+0.12), I(1.0), I(0.85),
        [{"t": num, "size": 30, "color": TEAL, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, I(x+1.15), I(y+0.16), I(cw-1.3), I(0.85),
        [{"t": h, "size": 14.5, "color": INK, "bold": True, "space_after": 2},
         {"t": b, "size": 10.5, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# SLIDE 3 — DEFINIÇÃO
# ======================================================================
s = blank(prs)
header(s, "Fundamentos", "Definição e conceitos", 3, TOTAL)
txt(s, I(0.86), I(1.75), I(7.1), I(1.3),
    [{"t": [{"t": "Invasão do espaço articular por microrganismos ", "size": 16, "color": TXT},
            {"t": "(bactérias, fungos ou micobactérias)", "size": 16, "color": TXT, "italic": True},
            {"t": ", com resposta inflamatória intensa que leva à destruição da cartilagem em horas a dias.",
             "size": 16, "color": TXT}],
      "line_spacing": 1.15}])
# emergência
box = rect(s, I(0.86), I(3.15), I(7.1), I(1.0), RGBColor_soft(ALERT), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.06)
rect(s, I(0.86), I(3.15), I(0.11), I(1.0), ALERT)
txt(s, I(1.15), I(3.28), I(6.7), I(0.8),
    [{"t": "EMERGÊNCIA ORTOPÉDICA", "size": 13, "color": ALERT, "bold": True, "space_after": 3},
     {"t": "A cartilagem hialina pode ser irreversivelmente lesada já nas primeiras 24–48 h. Diagnóstico e drenagem são tempo-dependentes.",
      "size": 12, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
# termos-chave à direita
tx = 8.35
txt(s, I(tx), I(1.7), I(4.3), I(0.4), [{"t": "TERMOS RELACIONADOS", "size": 12, "color": PRIMARY, "bold": True}])
terms = [
    ("Artrite séptica nativa", "Articulação sem implante"),
    ("Artrite gonocócica", "N. gonorrhoeae disseminada"),
    ("Infecção de prótese (PJI)", "Articulação com implante"),
    ("Pioartrose", "Sinônimo clássico"),
]
yy = 2.1
for h, b in terms:
    rect(s, I(tx), I(yy+0.06), I(0.12), I(0.5), TEAL)
    txt(s, I(tx+0.28), I(yy), I(4.0), I(0.6),
        [{"t": h, "size": 13, "color": INK, "bold": True, "space_after": 1},
         {"t": b, "size": 10.5, "color": MUTE}])
    yy += 0.78
citation(s, "Mathews CJ et al. Bacterial septic arthritis in adults. Lancet 2010;375:846–855.")


# ======================================================================
# SLIDE 4 — EPIDEMIOLOGIA
# ======================================================================
s = blank(prs)
header(s, "Fundamentos", "Epidemiologia", 4, TOTAL)
img(s, "epi_incidencia.png", 0.75, 1.95, w=7.2)
txt(s, I(0.86), I(1.72), I(7.2), I(0.35),
    [{"t": "Incidência estimada por grupo de risco", "size": 13, "color": TXT, "bold": True}])
# stats à direita
rx = 8.5
for (big, lab, col), yy in zip(
        [("2–10", "casos / 100.000 hab-ano\nna população geral", PRIMARY),
         ("↑ risco", "idade avançada, AR, diabetes,\nimunossupressão", AMBER),
         ("7–15%", "letalidade em adultos\n(até 50% se poliarticular)", ALERT)],
        [1.9, 3.42, 4.94]):
    card = rect(s, I(rx), I(yy), I(4.1), I(1.4), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.08)
    txt(s, I(rx+0.3), I(yy+0.2), I(1.7), I(1.0),
        [{"t": big, "size": 28, "color": col, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, I(rx+2.05), I(yy+0.2), I(1.9), I(1.0),
        [{"t": lab, "size": 11, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
citation(s, "Mathews CJ et al. Lancet 2010. · Kaandorp CJE et al. Arthritis Rheum 1997.")


# ======================================================================
# SLIDE 5 — CLASSIFICAÇÃO
# ======================================================================
s = blank(prs)
header(s, "Fundamentos", "Classificação", 5, TOTAL)
cols = [
    ("Segundo a articulação", TEAL, [("Nativa", "sem implante"),
                                     ("Protésica (PJI)", "com implante articular")]),
    ("Segundo o agente", PRIMARY, [("Gonocócica", "adulto jovem, sexualmente ativo"),
                                   ("Não gonocócica", "S. aureus e outros piogênicos")]),
    ("Segundo a evolução", AMBER, [("Aguda", "< 2 semanas, piogênica"),
                                   ("Crônica", "micobactérias, fungos, brucella")]),
]
cw = 3.85; gap = 0.28; x0 = 0.86; y0 = 2.0
for i, (title, col, rows) in enumerate(cols):
    x = x0 + i*(cw+gap)
    rect(s, I(x), I(y0), I(cw), I(0.7), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, I(x), I(y0+0.05), I(cw), I(0.6),
        [{"t": title, "size": 14, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE)
    yy = y0 + 0.95
    for h, b in rows:
        card = rect(s, I(x), I(yy), I(cw), I(1.35), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _set_corner(card, 0.07)
        txt(s, I(x+0.25), I(yy+0.2), I(cw-0.5), I(1.0),
            [{"t": h, "size": 15, "color": INK, "bold": True, "space_after": 4},
             {"t": b, "size": 11.5, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.5
citation(s, "Classificação com implicações microbiológicas, terapêuticas e prognósticas.")


# ======================================================================
# SLIDE 6 — MICROBIOLOGIA
# ======================================================================
s = blank(prs)
header(s, "Etiologia", "Microbiologia", 6, TOTAL, accent=ALERT)
img(s, "microbiologia.png", 0.75, 2.0, w=7.4)
txt(s, I(0.86), I(1.72), I(7.4), I(0.35),
    [{"t": "Agentes na artrite séptica nativa não gonocócica", "size": 13, "color": TXT, "bold": True}])
rx = 8.55
notes = [
    (ALERT, "S. aureus", "agente mais frequente; alta virulência e destruição rápida"),
    (PRIMARY, "Estreptococos", "S. pyogenes, S. agalactiae, S. pneumoniae"),
    (TEAL, "Gram-negativos", "idosos, UDIV, imunossuprimidos, ITU"),
    (AMBER, "Situações especiais", "gonococo (jovens), Salmonella (anemia falciforme)"),
]
yy = 1.95
for col, h, b in notes:
    rect(s, I(rx), I(yy+0.05), I(0.12), I(0.85), col)
    txt(s, I(rx+0.28), I(yy), I(4.1), I(1.0),
        [{"t": h, "size": 13.5, "color": INK, "bold": True, "space_after": 2},
         {"t": b, "size": 10.5, "color": TXT}])
    yy += 1.08
citation(s, "Distribuição aproximada a partir de séries de artrite séptica nativa (Mathews 2010; Ross 2017).")


# ======================================================================
# SLIDE 7 — FISIOPATOLOGIA (diagrama de vias)
# ======================================================================
s = blank(prs)
header(s, "Fisiopatologia", "Vias de infecção e lesão articular", 7, TOTAL, accent=ALERT)
# três vias -> articulação -> cascata
vias = [("Hematogênica", "a mais comum;\nfoco à distância", PRIMARY, 1.55),
        ("Inoculação direta", "trauma, punção,\nartroscopia, infiltração", TEAL, 3.05),
        ("Contiguidade", "osteomielite,\ncelulite adjacente", AMBER, 4.55)]
for h, b, col, yy in vias:
    node(s, I(0.86), I(yy), I(2.9), I(1.15), h + "\n" + b, col, size=12.5)
# articulação central
node(s, I(4.35), I(2.7), I(2.5), I(1.4), "Membrana\nsinovial\ninfectada", INK, size=15)
for _, _, _, yy in vias:
    connect(s, I(3.76), I(yy+0.57), I(4.35), I(3.4), color=MUTE)
# cascata à direita
casc = [("Resposta imune e enzimas", "PMN liberam proteases e citocinas"),
        ("Aumento da pressão intra-articular", "compromete perfusão da cartilagem"),
        ("Degradação da cartilagem", "condrólise em 24–48 h"),
        ("Dano articular irreversível", "artrose secundária, anquilose")]
connect(s, I(6.9), I(3.4), I(7.6), I(3.4), color=INK, width=2.2)
yy = 1.55
for i, (h, b) in enumerate(casc):
    col = [PRIMARY, AMBER, ALERT, INK][i]
    card = rect(s, I(7.75), I(yy), I(4.85), I(0.98), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.09)
    rect(s, I(7.75), I(yy), I(0.1), I(0.98), col)
    txt(s, I(8.0), I(yy+0.13), I(4.5), I(0.75),
        [{"t": f"{i+1}. {h}", "size": 12.5, "color": INK, "bold": True, "space_after": 1},
         {"t": b, "size": 10.5, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
    yy += 1.12
citation(s, "A via hematogênica responde pela maioria dos casos; a cartilagem é avascular e vulnerável.")


# ======================================================================
# SLIDE 8 — FATORES DE RISCO
# ======================================================================
s = blank(prs)
header(s, "Quadro clínico", "Fatores de risco", 8, TOTAL)
groups = [
    ("Do hospedeiro", PRIMARY, ["Idade > 80 anos", "Diabetes mellitus", "Imunossupressão / corticoide",
                                "Doença hepática ou renal crônica"]),
    ("Articulares", TEAL, ["Artrite reumatoide", "Osteoartrose avançada",
                           "Prótese articular", "Artrocentese / infiltração prévia"]),
    ("Comportamentais / cutâneos", AMBER, ["Uso de drogas injetáveis", "Bacteriemia / endocardite",
                                           "Úlceras e infecções de pele", "Cateteres e acessos vasculares"]),
]
cw = 3.85; gap = 0.28; x0 = 0.86; y0 = 1.95
for i, (title, col, items) in enumerate(groups):
    x = x0 + i*(cw+gap)
    card = rect(s, I(x), I(y0), I(cw), I(4.05), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.05)
    rect(s, I(x), I(y0), I(cw), I(0.62), col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, I(x), I(y0+0.02), I(cw), I(0.58),
        [{"t": title, "size": 13.5, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE)
    yy = y0 + 0.85
    for it in items:
        rect(s, I(x+0.3), I(yy+0.11), I(0.11), I(0.11), col, shape=MSO_SHAPE.OVAL)
        txt(s, I(x+0.55), I(yy), I(cw-0.75), I(0.55),
            [{"t": it, "size": 12, "color": TXT}])
        yy += 0.78
citation(s, "S. aureus tem tropismo aumentado por articulações reumatoides e protésicas.")


# ======================================================================
# SLIDE 9 — APRESENTAÇÃO CLÍNICA
# ======================================================================
s = blank(prs)
header(s, "Quadro clínico", "Apresentação clínica", 9, TOTAL)
# tríade
txt(s, I(0.86), I(1.75), I(6.5), I(0.35), [{"t": "TRÍADE CLÁSSICA", "size": 12.5, "color": ALERT, "bold": True}])
triad = [("Dor", "intensa, ao repouso e à mobilização"),
         ("Edema / derrame", "articulação quente e distendida"),
         ("Impotência funcional", "recusa ativa e passiva ao movimento")]
yy = 2.15
for h, b in triad:
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, I(0.86), I(yy), I(0.55), I(0.55))
    o.fill.solid(); o.fill.fore_color.rgb = ALERT; o.line.fill.background(); o.shadow.inherit = False
    txt(s, I(1.6), I(yy), I(5.6), I(0.6),
        [{"t": h, "size": 15.5, "color": INK, "bold": True, "space_after": 1},
         {"t": b, "size": 11.5, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.85
# pistas semiológicas
box = rect(s, I(0.86), I(4.95), I(6.4), I(1.15), RGBColor_soft(AMBER), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.06)
txt(s, I(1.1), I(5.08), I(6.0), I(0.95),
    [{"t": "⚠  ATENÇÃO", "size": 12, "color": AMBER, "bold": True, "space_after": 3},
     {"t": "Idosos, imunossuprimidos e portadores de prótese podem cursar sem febre. Monoartrite aguda = artrite séptica até prova em contrário.",
      "size": 11.5, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
# sistêmico à direita
rx = 7.6
card = rect(s, I(rx), I(1.95), I(5.0), I(4.15), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(card, 0.05)
txt(s, I(rx+0.3), I(2.15), I(4.5), I(0.4), [{"t": "SINAIS SISTÊMICOS E PADRÃO", "size": 12, "color": PRIMARY, "bold": True}])
sys = ["Febre (≈ 60%) e calafrios", "Mal-estar e taquicardia",
       "Monoartrite em > 80% dos casos", "Joelho e quadril são os mais afetados",
       "Poliarticular: pior prognóstico e maior letalidade"]
yy = 2.65
for it in sys:
    rect(s, I(rx+0.32), I(yy+0.09), I(0.1), I(0.1), PRIMARY, shape=MSO_SHAPE.OVAL)
    txt(s, I(rx+0.58), I(yy), I(4.2), I(0.6), [{"t": it, "size": 12, "color": TXT}])
    yy += 0.68
citation(s, "A ausência de febre não exclui o diagnóstico.")


# ======================================================================
# SLIDE 10 — ARTICULAÇÕES ACOMETIDAS
# ======================================================================
s = blank(prs)
header(s, "Quadro clínico", "Articulações mais acometidas", 10, TOTAL)
img(s, "articulacoes.png", 0.75, 2.05, w=7.4)
txt(s, I(0.86), I(1.75), I(7.4), I(0.35),
    [{"t": "Distribuição em adultos (articulação nativa)", "size": 13, "color": TXT, "bold": True}])
rx = 8.5
card = rect(s, I(rx), I(2.0), I(4.15), I(3.9), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(card, 0.05)
txt(s, I(rx+0.3), I(2.2), I(3.6), I(0.4), [{"t": "PONTOS PRÁTICOS", "size": 12, "color": PRIMARY, "bold": True}])
pts = ["Joelho é a articulação mais acometida no adulto",
       "Quadril predomina em lactentes e crianças pequenas",
       "Esterno-clavicular e sacroilíaca: pensar em UDIV",
       "Acometimento axial é raro e sugere agentes específicos"]
yy = 2.7
for it in pts:
    rect(s, I(rx+0.32), I(yy+0.09), I(0.1), I(0.1), TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, I(rx+0.58), I(yy), I(3.4), I(0.75), [{"t": it, "size": 11.5, "color": TXT}])
    yy += 0.8
citation(s, "UDIV: usuário de drogas injetáveis. Localizações atípicas exigem alta suspeição.")


# ======================================================================
# SLIDE 11 — DIAGNÓSTICO VISÃO GERAL
# ======================================================================
s = blank(prs)
header(s, "Diagnóstico", "Abordagem diagnóstica", 11, TOTAL, accent=TEAL)
steps = [
    ("1", "Suspeição clínica", "Monoartrite aguda + fatores de risco", PRIMARY),
    ("2", "Artrocentese", "ANTES do antibiótico, sempre que possível", ALERT),
    ("3", "Líquido sinovial", "Celularidade, Gram, cultura, cristais", TEAL),
    ("4", "Hemoculturas e exames", "PCR, VHS, hemograma, 2 hemoculturas", AMBER),
    ("5", "Imagem", "US, radiografia; RM/TC conforme sítio", PRIMARY),
]
cw = 2.28; gap = 0.14; x0 = 0.86; y = 2.35
for i, (n, h, b, col) in enumerate(steps):
    x = x0 + i*(cw+gap)
    card = rect(s, I(x), I(y), I(cw), I(2.5), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.07)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, I(x+cw/2-0.33), I(y+0.25), I(0.66), I(0.66))
    o.fill.solid(); o.fill.fore_color.rgb = col; o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    txt(s, I(x+0.15), I(y+1.05), I(cw-0.3), I(1.35),
        [{"t": h, "size": 13, "color": INK, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 4},
         {"t": b, "size": 10.5, "color": TXT, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.TOP)
    if i < len(steps)-1:
        connect(s, I(x+cw+0.005), I(y+0.58), I(x+cw+gap), I(y+0.58), color=MUTE)
box = rect(s, I(0.86), I(5.35), I(11.75), I(0.85), RGBColor_soft(ALERT), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.08)
txt(s, I(1.1), I(5.48), I(11.3), I(0.6),
    [{"t": [{"t": "Regra de ouro:  ", "size": 13, "color": ALERT, "bold": True},
            {"t": "toda monoartrite aguda deve ser puncionada. O antibiótico não deve atrasar a artrocentese diagnóstica.",
             "size": 12.5, "color": TXT}]}], anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# SLIDE 12 — MARCADORES SÉRICOS
# ======================================================================
s = blank(prs)
header(s, "Diagnóstico", "Marcadores séricos", 12, TOTAL, accent=TEAL)
rows = [
    ["Marcador", "Comportamento", "Utilidade clínica"],
    ["Leucograma", "Frequentemente ↑, mas pode ser normal", "Baixa especificidade"],
    ["VHS", "Elevada (> 30 mm/h)", "Sensível; útil no seguimento"],
    ["PCR", "Elevada; queda indica resposta", "Melhor para monitorar tratamento"],
    ["Procalcitonina", "Pode auxiliar se sepse", "Especificidade moderada"],
    ["Hemoculturas", "Positivas em até 50%", "Colher 2 pares antes do ATB"],
]
table(s, 0.86, 2.0, 7.5, rows, [1.7, 3.0, 2.8], font=11.5, row_h=0.62, head_h=0.5)
rx = 8.7
card = rect(s, I(rx), I(2.0), I(3.9), I(3.9), RGBColor_soft(PRIMARY), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(card, 0.05)
txt(s, I(rx+0.3), I(2.25), I(3.3), I(0.4), [{"t": "MENSAGEM", "size": 12, "color": PRIMARY, "bold": True}])
txt(s, I(rx+0.3), I(2.75), I(3.35), I(3.0),
    [{"t": "Nenhum marcador sérico confirma ou exclui isoladamente a artrite séptica.",
      "size": 13.5, "color": INK, "bold": True, "space_after": 8},
     {"t": "VHS e PCR normais reduzem a probabilidade, mas a decisão depende do líquido sinovial.",
      "size": 12, "color": TXT, "space_after": 8},
     {"t": "PCR seriada é o melhor parâmetro para acompanhar a resposta terapêutica.",
      "size": 12, "color": TXT}])
citation(s, "Carpenter CR et al. Acad Emerg Med 2011 (evidência sobre marcadores).")


# ======================================================================
# SLIDE 13 — LÍQUIDO SINOVIAL
# ======================================================================
s = blank(prs)
header(s, "Diagnóstico", "Artrocentese e líquido sinovial", 13, TOTAL, accent=ALERT)
img(s, "liquido_sinovial.png", 0.7, 1.95, w=6.5)
txt(s, I(0.86), I(1.72), I(7.0), I(0.35),
    [{"t": "Leucócitos por categoria de líquido articular", "size": 13, "color": TXT, "bold": True}])
rows = [
    ["Parâmetro", "Sugestivo de artrite séptica"],
    ["Aspecto", "Turvo / purulento"],
    ["Leucócitos", "> 50.000/mm³ (frequente > 100.000)"],
    ["Polimorfonucleares", "> 75–90%"],
    ["Gram", "Positivo em 30–50%"],
    ["Cultura", "Padrão-ouro (positiva 60–90%)"],
    ["Cristais", "Não excluem infecção concomitante"],
]
table(s, 7.95, 2.05, 4.7, rows, [1.85, 2.85], font=10.5, row_h=0.5, head_h=0.42)
# razões de verossimilhança (JAMA 2007)
txt(s, I(7.95), I(5.62), I(4.7), I(0.3),
    [{"t": "RAZÃO DE VEROSSIMILHANÇA (leucócitos)", "size": 10, "color": ALERT, "bold": True}])
lrs = [(">50.000", "LR 7,7"), (">100.000", "LR 28"), ("PMN >90%", "LR 3,4")]
for k, (lab, val) in enumerate(lrs):
    x = 7.95 + k*1.6
    chip = rect(s, I(x), I(5.92), I(1.45), I(0.62), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(chip, 0.14)
    txt(s, I(x), I(5.99), I(1.45), I(0.5),
        [{"t": val, "size": 13, "color": INK, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 0},
         {"t": lab, "size": 8.5, "color": MUTE, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
box = rect(s, I(0.86), I(5.32), I(6.5), I(1.08), RGBColor_soft(ALERT), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.06)
txt(s, I(1.1), I(5.42), I(6.1), I(0.9),
    [{"t": "Nenhum ponto de corte é absoluto", "size": 12.5, "color": ALERT, "bold": True, "space_after": 3},
     {"t": "Contagens < 50.000 não excluem infecção, sobretudo em imunossuprimidos e próteses. Sempre enviar cultura.",
      "size": 11.5, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
citation(s, "Margaretten ME et al. Does this adult patient have septic arthritis? JAMA 2007;297:1478–88.")


# ======================================================================
# SLIDE 14 — IMAGEM
# ======================================================================
s = blank(prs)
header(s, "Diagnóstico", "Métodos de imagem", 14, TOTAL, accent=TEAL)
mods = [
    ("Radiografia", PRIMARY, ["Normal no início", "Aumento de partes moles",
                              "Tardio: redução do espaço, erosões", "Exclui outras causas"]),
    ("Ultrassonografia", TEAL, ["Detecta derrame precoce", "Guia a artrocentese",
                                "Sem radiação, à beira-leito", "Método inicial de escolha"]),
    ("Ressonância / TC", AMBER, ["Alta sensibilidade para derrame", "Avalia osteomielite associada",
                                 "Abscessos e articulações profundas", "Sacroilíaca, quadril, ombro"]),
]
cw = 3.85; gap = 0.28; x0 = 0.86; y0 = 1.95
for i, (title, col, items) in enumerate(mods):
    x = x0 + i*(cw+gap)
    card = rect(s, I(x), I(y0), I(cw), I(4.05), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.05)
    rect(s, I(x), I(y0), I(cw), I(0.62), col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, I(x), I(y0+0.02), I(cw), I(0.58),
        [{"t": title, "size": 14, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE)
    yy = y0 + 0.9
    for it in items:
        rect(s, I(x+0.3), I(yy+0.1), I(0.11), I(0.11), col, shape=MSO_SHAPE.OVAL)
        txt(s, I(x+0.55), I(yy), I(cw-0.75), I(0.6), [{"t": it, "size": 11.5, "color": TXT}])
        yy += 0.78
citation(s, "A imagem apoia, mas nunca substitui a análise do líquido sinovial.")


# ======================================================================
# SLIDE 15 — ALGORITMO DIAGNÓSTICO DO ADULTO (flowchart)
# ======================================================================
s = blank(prs)
header(s, "Diagnóstico", "Algoritmo diagnóstico — adulto", 15, TOTAL, accent=PRIMARY)
node(s, I(4.9), I(1.5), I(3.5), I(0.66), "Monoartrite aguda / suspeita clínica", INK, size=13)
node(s, I(4.9), I(2.34), I(3.5), I(0.66), "Artrocentese (antes do ATB)", ALERT, size=13)
connect(s, I(6.65), I(2.16), I(6.65), I(2.34), color=MUTE)
node(s, I(4.9), I(3.18), I(3.5), I(0.66), "Análise do líquido sinovial", PRIMARY, size=13)
connect(s, I(6.65), I(3.0), I(6.65), I(3.18), color=MUTE)
# ramo positivo / negativo
node(s, I(1.15), I(4.35), I(3.6), I(0.98),
     "Purulento · Leucócitos > 50.000\nPMN > 90% · Gram/cultura +", ALERT, size=11.5)
node(s, I(8.55), I(4.35), I(3.6), I(0.98),
     "Baixa celularidade\nGram negativo, cristais +", GREEN, size=11.5)
connect(s, I(5.7), I(3.84), I(2.95), I(4.35), color=ALERT)
connect(s, I(7.6), I(3.84), I(10.35), I(4.35), color=GREEN)
edge_label(s, I(3.0), I(4.02), I(2.3), "compatível com séptica", ALERT)
edge_label(s, I(8.0), I(4.02), I(2.3), "outra hipótese", GREEN)
node(s, I(1.15), I(5.55), I(3.6), I(0.82),
     "Iniciar ATB empírico +\ndrenagem articular", ALERT, size=12)
connect(s, I(2.95), I(5.33), I(2.95), I(5.55), color=ALERT)
node(s, I(8.55), I(5.55), I(3.6), I(0.82),
     "Investigar gota, AR,\nreativa; reavaliar", GREEN, size=12)
connect(s, I(10.35), I(5.33), I(10.35), I(5.55), color=GREEN)
citation(s, "Se alta suspeita clínica, tratar mesmo com Gram negativo, aguardando cultura.")


# ======================================================================
# SLIDE 16 — ARTRITE SÉPTICA PEDIÁTRICA
# ======================================================================
s = blank(prs)
header(s, "Pediatria", "Artrite séptica na criança", 16, TOTAL, accent=AMBER)
left = [("Epidemiologia", "Mais comum < 3 anos; quadril e joelho"),
        ("Agentes", "S. aureus; Kingella kingae em lactentes; considerar S. do grupo B no neonato"),
        ("Quadril", "Emergência: risco de necrose avascular e luxação")]
yy = 1.95
for h, b in left:
    rect(s, I(0.86), I(yy+0.05), I(0.12), I(0.9), AMBER)
    txt(s, I(1.12), I(yy), I(5.8), I(1.0),
        [{"t": h, "size": 15, "color": INK, "bold": True, "space_after": 2},
         {"t": b, "size": 12, "color": TXT}])
    yy += 1.15
box = rect(s, I(0.86), I(5.5), I(6.0), I(0.9), RGBColor_soft(ALERT), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.07)
txt(s, I(1.1), I(5.6), I(5.6), I(0.72),
    [{"t": "Diferenciar de sinovite transitória do quadril — principal diagnóstico diferencial.",
      "size": 12, "color": TXT, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
# painel direito: apresentação
rx = 7.35
card = rect(s, I(rx), I(1.95), I(5.25), I(4.45), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(card, 0.04)
txt(s, I(rx+0.32), I(2.15), I(4.6), I(0.4), [{"t": "SINAIS NA CRIANÇA", "size": 12.5, "color": AMBER, "bold": True}])
kids = ["Recusa a apoiar / claudicação", "Pseudoparalisia do membro",
        "Postura antálgica (quadril em flexão e rotação externa)",
        "Irritabilidade e febre", "Dor à mobilização passiva"]
yy = 2.65
for it in kids:
    rect(s, I(rx+0.34), I(yy+0.09), I(0.11), I(0.11), AMBER, shape=MSO_SHAPE.OVAL)
    txt(s, I(rx+0.6), I(yy), I(4.4), I(0.6), [{"t": it, "size": 12.5, "color": TXT}])
    yy += 0.72
citation(s, "Kingella kingae exige meios de cultura específicos (frascos de hemocultura / PCR).")


# ======================================================================
# SLIDE 17 — CRITÉRIOS DE KOCHER
# ======================================================================
s = blank(prs)
header(s, "Pediatria", "Critérios de Kocher", 17, TOTAL, accent=AMBER)
txt(s, I(0.86), I(1.72), I(6.0), I(0.35),
    [{"t": "Diferenciam artrite séptica × sinovite transitória do quadril", "size": 12.5, "color": TXT, "bold": True}])
crit = ["Febre > 38,5 °C", "Recusa a apoiar o membro (não deambula)",
        "VHS > 40 mm/h", "Leucócitos séricos > 12.000/mm³"]
yy = 2.2
for i, c in enumerate(crit):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, I(0.86), I(yy), I(0.5), I(0.5))
    o.fill.solid(); o.fill.fore_color.rgb = AMBER; o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1); r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    txt(s, I(1.55), I(yy), I(5.3), I(0.55), [{"t": c, "size": 13.5, "color": INK, "bold": True}],
        anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.72
box = rect(s, I(0.86), I(5.2), I(6.0), I(1.15), RGBColor_soft(PRIMARY), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.06)
txt(s, I(1.1), I(5.31), I(5.6), I(0.95),
    [{"t": "Caird (2006): acrescenta PCR > 20 mg/L (modelo de 5 fatores).", "size": 11.5, "color": INK, "bold": True, "space_after": 3},
     {"t": "3 fatores ≈ 83%  ·  4 ≈ 93%  ·  5 ≈ 98% de probabilidade.", "size": 11.5, "color": TXT}],
    anchor=MSO_ANCHOR.MIDDLE)
img(s, "kocher.png", 7.15, 2.0, w=5.5)
txt(s, I(7.3), I(6.05), I(5.3), I(0.35),
    [{"t": "Probabilidade de artrite séptica por nº de critérios", "size": 11, "color": MUTE, "italic": True}])
citation(s, "Kocher MS et al. J Bone Joint Surg Am 1999. · Caird MS et al. JBJS Am 2006.", y=6.62)


# ======================================================================
# SLIDE 18 — ALGORITMO PEDIÁTRICO
# ======================================================================
s = blank(prs)
header(s, "Pediatria", "Algoritmo — quadril doloroso na criança", 18, TOTAL, accent=AMBER)
node(s, I(4.85), I(1.55), I(3.6), I(0.7), "Criança com quadril doloroso\ne recusa a apoiar", INK, size=12)
node(s, I(4.85), I(2.5), I(3.6), I(0.7), "Aplicar critérios de Kocher\n+ PCR (Caird)", AMBER, size=12)
connect(s, I(6.65), I(2.25), I(6.65), I(2.5), color=MUTE)
# três ramos
node(s, I(0.9), I(3.7), I(3.5), I(0.85), "0–1 critério\nBaixa probabilidade", GREEN, size=12)
node(s, I(4.85), I(3.7), I(3.6), I(0.85), "2 critérios\nProbabilidade intermediária", AMBER, size=12)
node(s, I(8.9), I(3.7), I(3.5), I(0.85), "3–4 critérios\nAlta probabilidade", ALERT, size=12)
connect(s, I(6.0), I(3.2), I(2.65), I(3.7), color=GREEN)
connect(s, I(6.65), I(3.2), I(6.65), I(3.7), color=AMBER)
connect(s, I(7.3), I(3.2), I(10.65), I(3.7), color=ALERT)
node(s, I(0.9), I(4.95), I(3.5), I(0.95), "Observação e\nreavaliação clínica", GREEN, size=12)
node(s, I(4.85), I(4.95), I(3.6), I(0.95), "USG + artrocentese\nguiada / RM", AMBER, size=12)
node(s, I(8.9), I(4.95), I(3.5), I(0.95), "Artrocentese urgente\n+ drenagem cirúrgica", ALERT, size=12)
for x in [2.65, 6.65, 10.65]:
    connect(s, I(x), I(4.55), I(x), I(4.95), color=MUTE)
citation(s, "A conduta cirúrgica no quadril não deve ser postergada quando a suspeição é alta.")


# ======================================================================
# SLIDE 19 — DIAGNÓSTICO DIFERENCIAL
# ======================================================================
s = blank(prs)
header(s, "Diagnóstico", "Diagnóstico diferencial", 19, TOTAL, accent=TEAL)
diffs = [
    ("Artrite por cristais", "Gota e pseudogota; cristais no líquido — podem coexistir"),
    ("Artrite reativa", "Pós-infecciosa; oligoartrite, entesite, uveíte"),
    ("Artrite reumatoide / AIJ", "Surto agudo pode simular infecção"),
    ("Hemartrose", "Trauma, anticoagulação, hemofilia"),
    ("Sinovite transitória", "Criança; autolimitada, sem toxemia"),
    ("Bursite / celulite", "Processo periarticular, sem derrame articular"),
]
cw = 5.75; gap = 0.35; x0 = 0.86; y0 = 1.95; ch = 1.2
for k, (h, b) in enumerate(diffs):
    col = k % 2; rowi = k // 2
    x = x0 + col*(cw+gap); y = y0 + rowi*(ch+0.18)
    card = rect(s, I(x), I(y), I(cw), I(ch), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.06)
    rect(s, I(x), I(y), I(0.1), I(ch), TEAL)
    txt(s, I(x+0.3), I(y+0.16), I(cw-0.55), I(0.9),
        [{"t": h, "size": 14, "color": INK, "bold": True, "space_after": 3},
         {"t": b, "size": 11.5, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
citation(s, "A presença de cristais NÃO exclui infecção — enviar sempre cultura.")


# ======================================================================
# SLIDE 20 — PRINCÍPIOS DO TRATAMENTO
# ======================================================================
s = blank(prs)
header(s, "Tratamento", "Princípios do tratamento", 20, TOTAL, accent=ALERT)
pillars = [
    ("Drenagem", ALERT, "Esvaziar e descomprimir a articulação: punção, artroscopia ou artrotomia"),
    ("Antibioticoterapia", PRIMARY, "Empírica precoce após culturas; ajuste dirigido pelo antibiograma"),
    ("Suporte e reabilitação", TEAL, "Analgesia, mobilização precoce, fisioterapia e controle de comorbidades"),
]
cw = 3.85; gap = 0.28; x0 = 0.86; y0 = 2.05
for i, (h, col, b) in enumerate(pillars):
    x = x0 + i*(cw+gap)
    card = rect(s, I(x), I(y0), I(cw), I(3.4), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.05)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, I(x+cw/2-0.5), I(y0+0.35), I(1.0), I(1.0))
    o.fill.solid(); o.fill.fore_color.rgb = col; o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1); r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    txt(s, I(x+0.2), I(y0+1.55), I(cw-0.4), I(0.5),
        [{"t": h, "size": 16, "color": INK, "bold": True, "align": PP_ALIGN.CENTER}])
    txt(s, I(x+0.35), I(y0+2.1), I(cw-0.7), I(1.2),
        [{"t": b, "size": 12, "color": TXT, "align": PP_ALIGN.CENTER}])
box = rect(s, I(0.86), I(5.75), I(11.75), I(0.72), RGBColor_soft(PRIMARY), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.1)
txt(s, I(1.1), I(5.83), I(11.3), I(0.56),
    [{"t": "Coletar culturas ANTES do antibiótico; não atrasar a drenagem à espera de exames.",
      "size": 12.5, "color": INK, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# SLIDE 21 — ANTIBIOTICOTERAPIA EMPÍRICA
# ======================================================================
s = blank(prs)
header(s, "Tratamento", "Antibioticoterapia empírica", 21, TOTAL, accent=PRIMARY)
txt(s, I(0.86), I(1.72), I(11), I(0.35),
    [{"t": "Guiada pelo Gram e pelo cenário clínico — ajustar sempre ao protocolo local", "size": 12.5, "color": TXT, "bold": True}])
rows = [
    ["Cenário / Gram", "Provável agente", "Esquema empírico usual"],
    ["Cocos Gram + / sem risco MRSA", "S. aureus, estreptococos", "Oxacilina ou cefazolina"],
    ["Risco de MRSA", "S. aureus resistente", "Vancomicina"],
    ["Bacilos Gram –", "Enterobactérias, Pseudomonas", "Cefalosporina 3ª/4ª geração"],
    ["Gram indisponível / grave", "Cobertura ampla", "Vancomicina + cefalosporina 3ª/4ª"],
    ["Adulto jovem sexualmente ativo", "N. gonorrhoeae", "Ceftriaxona"],
    ["Trauma / UDIV", "Incluir Pseudomonas", "Ativo anti-pseudomonas"],
]
table(s, 0.86, 2.25, 11.75, rows, [3.75, 3.5, 4.5], font=11.5, row_h=0.56, head_h=0.5)
citation(s, "Ajustar dose por peso e função renal. Descalonar após cultura/antibiograma.", y=6.5)


# ======================================================================
# SLIDE 22 — ANTIBIÓTICO DIRIGIDO E DURAÇÃO
# ======================================================================
s = blank(prs)
header(s, "Tratamento", "Terapia dirigida e duração", 22, TOTAL, accent=PRIMARY)
left = [
    ("Início intravenoso", "IV nas primeiras 2 semanas, conforme resposta clínica e laboratorial"),
    ("Transição para via oral", "Possível com boa evolução e agente sensível de boa biodisponibilidade"),
    ("Duração total", "Habitualmente 3–4 semanas (piogênicos); individualizar por agente e evolução"),
    ("Monitoramento", "PCR seriada, exame clínico e reavaliação da articulação"),
]
yy = 1.95
for h, b in left:
    card = rect(s, I(0.86), I(yy), I(6.4), I(1.02), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.06)
    rect(s, I(0.86), I(yy), I(0.1), I(1.02), PRIMARY)
    txt(s, I(1.12), I(yy+0.13), I(6.0), I(0.8),
        [{"t": h, "size": 13.5, "color": INK, "bold": True, "space_after": 2},
         {"t": b, "size": 11, "color": TXT}], anchor=MSO_ANCHOR.MIDDLE)
    yy += 1.12
rx = 7.65
card = rect(s, I(rx), I(1.95), I(4.95), I(4.5), RGBColor_soft(TEAL), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(card, 0.04)
txt(s, I(rx+0.32), I(2.2), I(4.3), I(0.4), [{"t": "AGENTES QUE EXIGEM MAIS TEMPO", "size": 12, "color": TEAL, "bold": True}])
dur = [("Gonococo", "7–14 dias"), ("Estreptococo / estafilococo", "3–4 semanas"),
       ("Gram-negativos", "3–4 semanas ou mais"), ("Micobactérias / fungos", "meses"),
       ("Prótese articular", "6–12 semanas + estratégia cirúrgica")]
yy = 2.7
for a, d in dur:
    txt(s, I(rx+0.32), I(yy), I(2.9), I(0.5), [{"t": a, "size": 12, "color": INK, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, I(rx+3.1), I(yy), I(1.7), I(0.5), [{"t": d, "size": 12, "color": TEAL, "bold": True, "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, I(rx+0.32), I(yy+0.5), I(4.3), Pt(0.8), LINE)
    yy += 0.62
citation(s, "Durações são orientativas; seguir protocolo institucional e evolução individual.")


# ======================================================================
# SLIDE 23 — DRENAGEM ARTICULAR
# ======================================================================
s = blank(prs)
header(s, "Tratamento", "Drenagem articular", 23, TOTAL, accent=ALERT)
opts = [
    ("Punção / aspiração", PRIMARY, ["Articulações superficiais acessíveis",
                                     "Pode ser repetida", "Falha → escalonar para cirurgia"]),
    ("Artroscopia", TEAL, ["Lavagem e desbridamento", "Menor morbidade",
                           "Preferida em joelho e ombro"]),
    ("Artrotomia aberta", ALERT, ["Quadril (crianças) e falha prévia",
                                  "Articulações profundas/loculadas", "Desbridamento amplo"]),
]
cw = 3.85; gap = 0.28; x0 = 0.86; y0 = 1.95
for i, (title, col, items) in enumerate(opts):
    x = x0 + i*(cw+gap)
    card = rect(s, I(x), I(y0), I(cw), I(3.6), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.05)
    rect(s, I(x), I(y0), I(cw), I(0.65), col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, I(x), I(y0+0.03), I(cw), I(0.6),
        [{"t": title, "size": 14, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    yy = y0 + 0.95
    for it in items:
        rect(s, I(x+0.3), I(yy+0.1), I(0.11), I(0.11), col, shape=MSO_SHAPE.OVAL)
        txt(s, I(x+0.55), I(yy), I(cw-0.75), I(0.7), [{"t": it, "size": 11.5, "color": TXT}])
        yy += 0.82
box = rect(s, I(0.86), I(5.85), I(11.75), I(0.7), RGBColor_soft(ALERT), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.1)
txt(s, I(1.1), I(5.92), I(11.3), I(0.56),
    [{"t": "Quadril séptico é indicação clássica de drenagem cirúrgica urgente, especialmente em crianças.",
      "size": 12.5, "color": INK, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# SLIDE 24 — ALGORITMO DE MANEJO CIRÚRGICO
# ======================================================================
s = blank(prs)
header(s, "Tratamento", "Algoritmo de manejo — articulação nativa", 24, TOTAL, accent=ALERT)
node(s, I(4.85), I(1.5), I(3.6), I(0.72), "Diagnóstico de artrite séptica", INK, size=13)
node(s, I(4.85), I(2.42), I(3.6), I(0.72), "ATB empírico + drenagem inicial", ALERT, size=12.5)
connect(s, I(6.65), I(2.22), I(6.65), I(2.42), color=MUTE)
node(s, I(4.85), I(3.34), I(3.6), I(0.72), "Reavaliar em 48–72 h", PRIMARY, size=13)
connect(s, I(6.65), I(3.14), I(6.65), I(3.34), color=MUTE)
node(s, I(1.2), I(4.5), I(3.6), I(1.0), "Melhora clínica e\nqueda da PCR", GREEN, size=12)
node(s, I(8.5), I(4.5), I(3.6), I(1.0), "Sem melhora /\nderrame persistente", ALERT, size=12)
connect(s, I(5.7), I(4.06), I(3.0), I(4.5), color=GREEN)
connect(s, I(7.6), I(4.06), I(10.3), I(4.5), color=ALERT)
node(s, I(1.2), I(5.6), I(3.6), I(0.8), "Manter ATB e\nreabilitação", GREEN, size=12)
node(s, I(8.5), I(5.6), I(3.6), I(0.8), "Nova drenagem /\ndesbridamento cirúrgico", ALERT, size=12)
connect(s, I(3.0), I(5.5), I(3.0), I(5.6), color=GREEN)
connect(s, I(10.3), I(5.5), I(10.3), I(5.6), color=ALERT)
citation(s, "Falha de resposta impõe reintervenção e reavaliação do agente e do foco.")


# ======================================================================
# SLIDE 25 — ARTRITE GONOCÓCICA
# ======================================================================
s = blank(prs)
header(s, "Cenários avançados", "Artrite gonocócica", 25, TOTAL, accent=TEAL)
txt(s, I(0.86), I(1.72), I(6.3), I(0.6),
    [{"t": "Infecção gonocócica disseminada — adulto jovem, sexualmente ativo",
      "size": 13, "color": TXT, "bold": True}])
# duas formas
forms = [("Forma artrite-dermatite", ["Poliartralgia migratória", "Tenossinovite", "Lesões cutâneas pustulosas", "Cultura articular geralmente negativa"]),
         ("Forma artrite purulenta", ["Mono/oligoartrite", "Líquido purulento", "Cultura pode ser positiva", "Menos manifestações sistêmicas"])]
yy = 2.4
for title, items in forms:
    card = rect(s, I(0.86), I(yy), I(6.3), I(1.85), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.05)
    rect(s, I(0.86), I(yy), I(0.1), I(1.85), TEAL)
    txt(s, I(1.12), I(yy+0.12), I(5.9), I(0.4), [{"t": title, "size": 13.5, "color": INK, "bold": True}])
    xx = 1.15; ry = yy+0.58
    for k, it in enumerate(items):
        rect(s, I(xx), I(ry+0.08), I(0.1), I(0.1), TEAL, shape=MSO_SHAPE.OVAL)
        txt(s, I(xx+0.24), I(ry), I(2.8), I(0.5), [{"t": it, "size": 10.5, "color": TXT}])
        ry += 0.55
        if k == 1:
            xx = 3.95; ry = yy+0.58
    yy += 2.05
# manejo à direita
rx = 7.55
card = rect(s, I(rx), I(1.95), I(5.05), I(4.5), RGBColor_soft(TEAL), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(card, 0.04)
txt(s, I(rx+0.32), I(2.2), I(4.4), I(0.4), [{"t": "DIAGNÓSTICO E TRATAMENTO", "size": 12.5, "color": TEAL, "bold": True}])
gon = ["NAAT/cultura de sítios genital, retal e faríngeo",
       "Hemoculturas e cultura do líquido sinovial",
       "Ceftriaxona é o tratamento de escolha",
       "Rastrear e tratar coinfecção por clamídia",
       "Boa resposta clínica; excelente prognóstico",
       "Notificar e tratar parceiros (IST)"]
yy = 2.7
for it in gon:
    rect(s, I(rx+0.34), I(yy+0.1), I(0.11), I(0.11), TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, I(rx+0.6), I(yy), I(4.2), I(0.6), [{"t": it, "size": 11.5, "color": TXT}])
    yy += 0.62
citation(s, "Melhor prognóstico articular que a artrite não gonocócica.")


# ======================================================================
# SLIDE 26 — INFECÇÃO DE PRÓTESE (PJI)
# ======================================================================
s = blank(prs)
header(s, "Cenários avançados", "Infecção de prótese articular (PJI)", 26, TOTAL, accent=ALERT)
# classificação temporal
txt(s, I(0.86), I(1.72), I(11), I(0.35),
    [{"t": "Classificação temporal (a partir da cirurgia)", "size": 12.5, "color": TXT, "bold": True}])
tl = [("Precoce", "< 3 meses", "S. aureus, Gram-negativos", PRIMARY),
      ("Tardia (retardada)", "3–12/24 meses", "Estafilococos coag.-neg., Cutibacterium", TEAL),
      ("Tardia hematogênica", "> 12–24 meses", "Disseminação de foco à distância", AMBER)]
cw = 3.85; gap = 0.28; x0 = 0.86; y0 = 2.15
for i, (h, t, ag, col) in enumerate(tl):
    x = x0 + i*(cw+gap)
    card = rect(s, I(x), I(y0), I(cw), I(1.75), CARDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_corner(card, 0.06)
    rect(s, I(x), I(y0), I(cw), I(0.55), col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, I(x), I(y0+0.02), I(cw), I(0.5), [{"t": h, "size": 13.5, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, I(x+0.2), I(y0+0.68), I(cw-0.4), I(1.0),
        [{"t": t, "size": 15, "color": INK, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 3},
         {"t": ag, "size": 10.5, "color": TXT, "align": PP_ALIGN.CENTER}])
# critérios diagnósticos
box = rect(s, I(0.86), I(4.25), I(11.75), I(2.0), RGBColor_soft(ALERT), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.04)
txt(s, I(1.15), I(4.42), I(11), I(0.4), [{"t": "CRITÉRIOS DIAGNÓSTICOS (EBJIS / ICM — confirmam PJI)", "size": 12.5, "color": ALERT, "bold": True}])
crit = ["Fístula/trajeto comunicante com a prótese ou pus visível",
        "≥ 2 culturas periprotéticas com o mesmo microrganismo",
        "Líquido sinovial > 3.000 leuc./mm³ ou PMN > 80% (PJI crônica)",
        "Histologia periprotética com inflamação aguda",
        "Biomarcadores sinoviais (ex.: alfa-defensina)"]
xx = 1.2; yy = 4.9; col = 0
for k, c in enumerate(crit):
    rect(s, I(xx), I(yy+0.08), I(0.11), I(0.11), ALERT, shape=MSO_SHAPE.OVAL)
    txt(s, I(xx+0.26), I(yy), I(5.3), I(0.5), [{"t": c, "size": 11, "color": TXT}])
    yy += 0.44
    if k == 2:
        xx = 6.9; yy = 4.9
citation(s, "MSIS/ICM 2018 e EBJIS 2021 definem PJI. Limiares sinoviais diferem da articulação nativa.")


# ======================================================================
# SLIDE 27 — PJI: ALGORITMO DE MANEJO
# ======================================================================
s = blank(prs)
header(s, "Cenários avançados", "PJI — estratégias de tratamento", 27, TOTAL, accent=ALERT)
node(s, I(4.85), I(1.5), I(3.6), I(0.72), "PJI confirmada", INK, size=13)
node(s, I(1.0), I(2.75), I(3.5), I(1.05),
     "Aguda / precoce\nImplante estável\nsem trajeto fistuloso", PRIMARY, size=11.5)
node(s, I(4.9), I(2.75), I(3.5), I(1.05),
     "Crônica\nou implante instável", TEAL, size=12)
node(s, I(8.8), I(2.75), I(3.5), I(1.05),
     "Comorbidade proibitiva\nou recusa cirúrgica", AMBER, size=11.5)
connect(s, I(5.9), I(2.22), I(2.75), I(2.75), color=MUTE)
connect(s, I(6.65), I(2.22), I(6.65), I(2.75), color=MUTE)
connect(s, I(7.4), I(2.22), I(10.55), I(2.75), color=MUTE)
node(s, I(1.0), I(4.25), I(3.5), I(1.0),
     "DAIR\ndesbridamento + troca de\ncomponentes móveis + ATB", PRIMARY, size=11)
node(s, I(4.9), I(4.25), I(3.5), I(1.0),
     "Revisão em 1 ou 2 tempos\n(espaçador com antibiótico)", TEAL, size=11.5)
node(s, I(8.8), I(4.25), I(3.5), I(1.0),
     "Terapia supressiva\ncrônica com antibiótico", AMBER, size=11.5)
for x in [2.75, 6.65, 10.55]:
    connect(s, I(x), I(3.8), I(x), I(4.25), color=MUTE)
box = rect(s, I(0.86), I(5.65), I(11.75), I(0.75), RGBColor_soft(PRIMARY), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.1)
txt(s, I(1.1), I(5.73), I(11.3), I(0.6),
    [{"t": "Decisão multidisciplinar (ortopedia + infectologia). Retenção do implante exige agudeza, estabilidade e agente tratável com biofilme (ex.: rifampicina para estafilococos).",
      "size": 11.5, "color": INK, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# SLIDE 28 — COMPLICAÇÕES E PROGNÓSTICO
# ======================================================================
s = blank(prs)
header(s, "Desfechos", "Complicações e prognóstico", 28, TOTAL, accent=ALERT)
img(s, "prognostico.png", 0.75, 2.2, w=6.0)
txt(s, I(0.86), I(1.9), I(6.6), I(0.35),
    [{"t": "Desfechos desfavoráveis (aprox.)", "size": 13, "color": TXT, "bold": True}])
rx = 7.7
txt(s, I(rx), I(1.9), I(5.0), I(0.4), [{"t": "FATORES DE PIOR PROGNÓSTICO", "size": 12.5, "color": ALERT, "bold": True}])
factors = ["Idade avançada e comorbidades", "Atraso no diagnóstico e na drenagem",
           "Acometimento poliarticular", "S. aureus, sobretudo MRSA",
           "Prótese articular e imunossupressão", "Doença articular prévia (ex.: AR)"]
yy = 2.4
for it in factors:
    rect(s, I(rx), I(yy+0.09), I(0.11), I(0.11), ALERT, shape=MSO_SHAPE.OVAL)
    txt(s, I(rx+0.26), I(yy), I(4.7), I(0.6), [{"t": it, "size": 12, "color": TXT}])
    yy += 0.6
box = rect(s, I(rx), I(6.05), I(5.0), I(0.5), RGBColor_soft(AMBER), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_set_corner(box, 0.12)
txt(s, I(rx+0.2), I(6.1), I(4.7), I(0.4),
    [{"t": "Diagnóstico e drenagem precoces = melhor desfecho funcional.", "size": 10.5, "color": INK, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
citation(s, "Sequelas: rigidez, artrose secundária, osteomielite e instabilidade articular.")


# ======================================================================
# SLIDE 29 — MENSAGENS-CHAVE
# ======================================================================
s = blank(prs)
rect(s, 0, 0, prs.slide_width, prs.slide_height, INK)
rect(s, 0, 0, I(0.35), prs.slide_height, TEAL)
txt(s, I(0.9), I(0.7), I(11.5), I(0.9),
    [{"t": "MENSAGENS-CHAVE", "size": 14, "color": TEAL, "bold": True, "space_after": 4},
     {"t": "10 pontos para levar para a prática", "size": 26, "color": WHITE, "bold": True}])
msgs = [
    "Monoartrite aguda é artrite séptica até prova em contrário.",
    "Puncionar sempre antes do antibiótico; enviar Gram, cultura e cristais.",
    "Nenhum ponto de corte de leucócitos exclui infecção isoladamente.",
    "S. aureus é o agente mais comum na forma não gonocócica.",
    "Joelho no adulto; quadril na criança — este último é emergência.",
    "Kocher + Caird orientam a criança com quadril doloroso.",
    "Tratamento = drenagem + antibiótico dirigido + reabilitação.",
    "Ceftriaxona para a artrite gonocócica; ótimo prognóstico.",
    "PJI exige critérios formais e decisão com a infectologia.",
    "Tempo é cartilagem: atraso piora o desfecho funcional.",
]
cw = 5.75; gap = 0.35; x0 = 0.9; y0 = 2.05; ch = 0.82
for k, m in enumerate(msgs):
    col = k % 2; rowi = k // 2
    x = x0 + col*(cw+gap); y = y0 + rowi*(ch+0.12)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y+0.16), I(0.5), I(0.5))
    o.fill.solid(); o.fill.fore_color.rgb = TEAL; o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(k+1); r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    txt(s, I(x+0.68), I(y), I(cw-0.75), I(0.82),
        [{"t": m, "size": 12, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# SLIDE 30 — REFERÊNCIAS
# ======================================================================
s = blank(prs)
header(s, "Referências", "Referências", 30, TOTAL, accent=PRIMARY)
refs = [
    "1. Margaretten ME, Kohlwes J, Moore D, Bent S. Does this adult patient have septic arthritis? JAMA. 2007;297(13):1478–88. PMID 17405973.",
    "2. Mathews CJ, Weston VC, Jones A, Field M, Coakley G. Bacterial septic arthritis in adults. Lancet. 2010;375(9717):846–55. PMID 20206778.",
    "3. Coakley G, et al. BSR & BHPR, BOA, RCGP and BSAC guidelines for management of the hot swollen joint in adults. Rheumatology (Oxford). 2006;45(8):1039–41. PMID 16829534.",
    "4. Kocher MS, Zurakowski D, Kasser JR. Differentiating septic arthritis from transient synovitis of the hip in children. J Bone Joint Surg Am. 1999;81(12):1662–70. PMID 10608376.",
    "5. Caird MS, et al. Factors distinguishing septic arthritis from transient synovitis of the hip in children: a prospective study. J Bone Joint Surg Am. 2006;88(6):1251–7. PMID 16757758.",
    "6. Luhmann SJ, et al. Differentiation between septic arthritis and transient synovitis of the hip in children. J Bone Joint Surg Am. 2004;86(5):956–62. PMID 15118038.",
    "7. Davis CM, Zamora RA. Surgical options and approaches for septic arthritis of the native hip and knee joint. J Arthroplasty. 2020;35(3S):S14–S18. PMID 32046824.",
    "8. Tverring J, et al. Septic Arthritis Score (SAS) for the adult native knee. BMC Infect Dis. 2025;25(1):926. PMID 40681984.",
    "9. McNally M, et al. The EBJIS definition of periprosthetic joint infection. Bone Joint J. 2021;103-B(1):18–25. PMID 33380199.",
    "10. Parvizi J, et al. The 2018 definition of periprosthetic hip and knee infection. J Arthroplasty. 2018;33(5):1309–14. PMID 29551303.",
    "11. Osmon DR, et al. Diagnosis and management of prosthetic joint infection: IDSA guidelines. Clin Infect Dis. 2013;56(1):e1–e25. PMID 23223583.",
    "12. Workowski KA, et al. Sexually transmitted infections treatment guidelines, 2021. MMWR Recomm Rep. 2021;70(4):1–187.",
]
col_w = 5.9; x0 = 0.86; y0 = 1.9
half = (len(refs)+1)//2
for k, r in enumerate(refs):
    col = 0 if k < half else 1
    idx_in = k if k < half else k - half
    x = x0 + col*(col_w+0.35)
    y = y0 + idx_in*0.85
    txt(s, I(x), I(y), I(col_w), I(0.82), [{"t": r, "size": 9.5, "color": TXT, "line_spacing": 1.02}])
citation(s, "Consultar diretrizes locais/institucionais atualizadas antes de decisões terapêuticas.", y=6.62)


# guardar
from pptx.dml.color import RGBColor  # noqa
out = os.path.join(os.path.dirname(__file__), "Artrite_Septica_Diego_Idarraga.pptx")
prs.save(out)
print("SAVED", out, "slides:", len(prs.slides._sldIdLst))
