from __future__ import annotations

from pathlib import Path
import re
import shlex
import subprocess

SAFE_PREFIXES = {
    'pwd', 'whoami', 'uname', 'date', 'ls', 'find', 'git', 'gh', 'node', 'npm', 'npx', 'python3', 'pip', 'pip3', 'ollama',
    'open', 'cat', 'head', 'tail', 'wc', 'which', 'env', 'ffmpeg', 'ffprobe', 'say', 'mdfind', 'cp', 'mv', 'mkdir', 'touch'
}

DEFAULT_POWERPOINT_TEMPLATE = Path('/Users/diegoalejandroidarragalopez/Desktop/Fraturas_por_Estresse_Uma_Visão_Abrangente_para_o_Ortopedista.pptx')

APP_ALIASES = {
    'excel': 'Microsoft Excel',
    'xlsx': 'Microsoft Excel',
    'spreadsheet': 'Microsoft Excel',
    'hoja de calculo': 'Microsoft Excel',
    'hoja de cálculo': 'Microsoft Excel',
    'microsoft excel': 'Microsoft Excel',
    'powerpoint': 'Microsoft PowerPoint',
    'presentacion': 'Microsoft PowerPoint',
    'presentación': 'Microsoft PowerPoint',
    'slides': 'Microsoft PowerPoint',
    'ppt': 'Microsoft PowerPoint',
    'pptx': 'Microsoft PowerPoint',
    'microsoft powerpoint': 'Microsoft PowerPoint',
    'word': 'Microsoft Word',
    'docx': 'Microsoft Word',
    'documento': 'Microsoft Word',
    'microsoft word': 'Microsoft Word',
    'brave': 'Brave Browser',
    'browser': 'Brave Browser',
    'telegram': 'Telegram',
    'whatsapp': 'WhatsApp',
    'finder': 'Finder',
    'preview': 'Preview',
    'keynote': 'Keynote',
    'numbers': 'Numbers',
}


def normalize_command(task: dict) -> str | None:
    inputs = task.get('inputs', {}) or {}
    for key in ['command', 'shell', 'cmd']:
        value = inputs.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    objective = (task.get('objective') or '').strip()
    if objective.startswith('!'):
        return objective[1:].strip()
    if objective:
        try:
            head = shlex.split(objective)[0]
        except Exception:
            head = ''
        if head in SAFE_PREFIXES:
            return objective
    return None


def run_safe_command(command: str, *, cwd: str | Path | None = None, timeout: int = 120) -> dict:
    head = shlex.split(command)[0] if command.strip() else ''
    if head not in SAFE_PREFIXES:
        return {
            'command': command,
            'allowed': False,
            'returncode': -1,
            'stdout': '',
            'stderr': f'Command not allowlisted for ops-runner v1: {head}',
        }

    completed = subprocess.run(
        command,
        cwd=str(cwd) if cwd else None,
        capture_output=True,
        text=True,
        timeout=timeout,
        shell=True,
        check=False,
    )
    return {
        'command': command,
        'allowed': True,
        'returncode': completed.returncode,
        'stdout': completed.stdout.strip(),
        'stderr': completed.stderr.strip(),
    }


def _resolve_app_name(value: str | None) -> str | None:
    if not value:
        return None
    key = value.strip().lower()
    return APP_ALIASES.get(key, value.strip())


def _default_office_output(kind: str) -> Path:
    desktop = Path.home() / 'Desktop'
    suffix = {'excel': '.xlsx', 'powerpoint': '.pptx', 'word': '.docx'}[kind]
    return desktop / f'orquesta-{kind}{suffix}'


def _run_subprocess(args: list[str], *, timeout: int = 60) -> dict:
    completed = subprocess.run(args, capture_output=True, text=True, timeout=timeout, check=False)
    return {
        'command': shlex.join(args),
        'returncode': completed.returncode,
        'stdout': completed.stdout.strip(),
        'stderr': completed.stderr.strip(),
    }


def run_applescript(script: str, *, timeout: int = 60) -> dict:
    return _run_subprocess(['osascript', '-e', script], timeout=timeout)


def create_excel_file(target: str | Path) -> dict:
    from openpyxl import Workbook

    path = Path(target).expanduser().resolve()
    path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = 'Sheet1'
    ws['A1'] = 'Orquesta'
    ws['A2'] = 'Listo para editar'
    wb.save(path)
    return {'path': str(path), 'kind': 'excel'}


def create_word_file(target: str | Path, text: str | None = None) -> dict:
    from docx import Document

    path = Path(target).expanduser().resolve()
    path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    doc.add_heading('Orquesta Document', level=1)
    if text:
        for paragraph in [item.strip() for item in text.split('|') if item.strip()]:
            doc.add_paragraph(paragraph)
    else:
        doc.add_paragraph('Documento creado por Orquesta.')
    doc.save(path)
    return {'path': str(path), 'kind': 'word'}


def create_powerpoint_file(target: str | Path, text: str | None = None) -> dict:
    from pptx import Presentation
    from pptx.util import Inches

    path = Path(target).expanduser().resolve()
    path.parent.mkdir(parents=True, exist_ok=True)
    template_used = DEFAULT_POWERPOINT_TEMPLATE if DEFAULT_POWERPOINT_TEMPLATE.exists() else None
    prs = Presentation(str(template_used)) if template_used else Presentation()
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    chunks = [item.strip() for item in (text or 'Presentacion creada por Orquesta').split('|') if item.strip()]
    title = chunks[0]
    bullets = chunks[1:] or ['Lista para editar']

    if getattr(slide.shapes, 'title', None):
        slide.shapes.title.text = title
    else:
        slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.5), Inches(0.8)).text_frame.text = title

    body_placeholder = None
    for shape in slide.placeholders:
        if getattr(shape, 'placeholder_format', None) and shape.placeholder_format.idx != 0:
            body_placeholder = shape
            break

    if body_placeholder and hasattr(body_placeholder, 'text_frame'):
        body = body_placeholder.text_frame
        body.clear()
        for idx, item in enumerate(bullets):
            if idx == 0:
                body.text = item
            else:
                body.add_paragraph().text = item
    else:
        body = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(4.5)).text_frame
        for idx, item in enumerate(bullets):
            if idx == 0:
                body.text = item
            else:
                body.add_paragraph().text = item
    prs.save(path)
    return {
        'path': str(path),
        'kind': 'powerpoint',
        'template_used': str(template_used) if template_used else None,
        'slide_count': len(prs.slides),
    }


def execute_app_action(action: str, *, app: str | None = None, target: str | None = None, timeout: int = 90) -> dict:
    resolved_app = _resolve_app_name(app)
    resolved_target = str(Path(target).expanduser()) if target and not str(target).startswith(('http://', 'https://')) else target

    if action == 'open_app' and resolved_app:
        return _run_subprocess(['open', '-a', resolved_app], timeout=timeout)
    if action == 'activate_app' and resolved_app:
        return run_applescript(f'tell application "{resolved_app}" to activate', timeout=timeout)
    if action == 'quit_app' and resolved_app:
        return run_applescript(f'tell application "{resolved_app}" to quit', timeout=timeout)
    if action == 'open_with_app' and resolved_app and resolved_target:
        return _run_subprocess(['open', '-a', resolved_app, resolved_target], timeout=timeout)
    if action == 'open_path' and resolved_target:
        return _run_subprocess(['open', resolved_target], timeout=timeout)
    if action == 'reveal_path' and resolved_target:
        return _run_subprocess(['open', '-R', resolved_target], timeout=timeout)
    if action == 'open_url' and resolved_target:
        return _run_subprocess(['open', resolved_target], timeout=timeout)
    if action == 'create_excel':
        created = create_excel_file(resolved_target or _default_office_output('excel'))
        opened = _run_subprocess(['open', '-a', 'Microsoft Excel', created['path']], timeout=timeout)
        return {**opened, 'created_path': created['path']}
    if action == 'create_word':
        created = create_word_file(resolved_target or _default_office_output('word'), text=app)
        opened = _run_subprocess(['open', '-a', 'Microsoft Word', created['path']], timeout=timeout)
        return {**opened, 'created_path': created['path']}
    if action == 'create_powerpoint':
        created = create_powerpoint_file(resolved_target or _default_office_output('powerpoint'), text=app)
        opened = _run_subprocess(['open', '-a', 'Microsoft PowerPoint', created['path']], timeout=timeout)
        return {**opened, 'created_path': created['path']}
    return {
        'command': f'{action} {resolved_app or ""} {resolved_target or ""}'.strip(),
        'returncode': -1,
        'stdout': '',
        'stderr': 'Unsupported or incomplete app action',
    }


def infer_app_request(task: dict) -> dict | None:
    inputs = task.get('inputs', {}) or {}
    if inputs.get('app_action'):
        return {
            'action': inputs.get('app_action'),
            'app': inputs.get('app_name'),
            'target': inputs.get('app_target'),
        }

    objective = (task.get('objective') or '').strip()
    if ';' in objective:
        objective = objective.split(';')[-1].strip()
    lowered = objective.lower()

    patterns = [
        (r'^(?:crea|crear)\s+excel(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_excel'),
        (r'^(?:crea|crear)\s+powerpoint(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_powerpoint'),
        (r'^(?:crea|crear)\s+word(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_word'),
        (r'^(?:crea|crear|haz|genera)\s+(?:un\s+)?(?:excel|xlsx|spreadsheet|hoja de calculo|hoja de cálculo)(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_excel'),
        (r'^(?:crea|crear|haz|genera)\s+(?:una\s+)?(?:powerpoint|presentacion(?:\s+powerpoint)?|presentación(?:\s+powerpoint)?|ppt|pptx|slides|diapositivas)(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_powerpoint'),
        (r'^(?:crea|crear|haz|genera)\s+(?:un\s+)?(?:word|docx|documento)(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_word'),
        (r'^(?:hazme|generame|genérame|creame|créame)\s+(?:un\s+)?(?:excel|xlsx|spreadsheet|hoja de calculo|hoja de cálculo)(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_excel'),
        (r'^(?:hazme|generame|genérame|creame|créame)\s+(?:una\s+)?(?:powerpoint|presentacion(?:\s+powerpoint)?|presentación(?:\s+powerpoint)?|ppt|pptx|slides|diapositivas)(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_powerpoint'),
        (r'^(?:hazme|generame|genérame|creame|créame)\s+(?:un\s+)?(?:word|docx|documento)(?:\s+en\s+(.+?))?(?:\s*:\s*(.+))?$', 'create_word'),
        (r'^(?:abre|open)\s+([a-z0-9 ._-]+?)\s+con\s+(.+)$', 'open_with_app'),
        (r'^(?:abre|open)\s+([a-z0-9 ._-]+?)$', 'open_app'),
        (r'^(?:activa|activate)\s+([a-z0-9 ._-]+?)$', 'activate_app'),
        (r'^(?:cierra|quit|cierra app)\s+([a-z0-9 ._-]+?)$', 'quit_app'),
        (r'^(?:abre archivo|open file)\s+(.+)$', 'open_path'),
        (r'^(?:muestra archivo|reveal file|revela)\s+(.+)$', 'reveal_path'),
        (r'^(?:abre carpeta|open folder)\s+(.+)$', 'open_path'),
        (r'^(?:muestra carpeta|reveal folder)\s+(.+)$', 'reveal_path'),
        (r'^(?:abre web|open url|abre url)\s+(.+)$', 'open_url'),
    ]

    for pattern, action in patterns:
        match = re.match(pattern, objective, flags=re.IGNORECASE)
        if not match:
            continue
        if action in {'create_excel', 'create_powerpoint', 'create_word'}:
            return {'action': action, 'app': (match.group(2) if match.lastindex and match.lastindex >= 2 else None), 'target': match.group(1) if match.group(1) else None}
        if action in {'open_app', 'activate_app', 'quit_app'}:
            return {'action': action, 'app': match.group(1).strip(), 'target': None}
        if action == 'open_with_app':
            return {'action': action, 'app': match.group(1).strip(), 'target': match.group(2).strip()}
        return {'action': action, 'app': None, 'target': match.group(1).strip()}

    if any(token in lowered for token in APP_ALIASES):
        for token in APP_ALIASES:
            if token in lowered:
                return {'action': 'open_app', 'app': token, 'target': None}
    return None
