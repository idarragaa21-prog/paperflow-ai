from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def main() -> None:
    out = Path(__file__).resolve().parents[1] / "app" / "templates" / "minimal_medical.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen (13.333 x 7.5)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "TÍTULO"
    subtitle.text = "Subtítulo / Autor / Institución"
    for p in title.text_frame.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0B, 0x2A, 0x4A)
    for p in subtitle.text_frame.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Content slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Título de slide"
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• Bullet principal"
    p.font.name = "Calibri"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
    p2 = tf.add_paragraph()
    p2.text = "• Bullet secundario"
    p2.font.name = "Calibri"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # Section header slide (reuse layout 5 if exists, else blank)
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Add a big centered title box
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "SECCIÓN"
    p.font.name = "Calibri"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0B, 0x2A, 0x4A)
    p.alignment = PP_ALIGN.CENTER

    # References slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Referencias"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "1. Referencia Vancouver..."
    p.font.name = "Calibri"
    p.font.size = Pt(14)

    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
