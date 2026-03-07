from __future__ import annotations

import json
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any
from uuid import UUID

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.storage import storage_manager


MAX_SLIDES = 50


def validate_outline(o: dict[str, Any]) -> list[str]:
    """Strict-ish validation to avoid generic decks.

    Returns list of human-readable errors.
    """

    errors: list[str] = []
    slides = o.get("slides", [])
    if not isinstance(slides, list):
        return ["outline.slides must be a list"]

    for idx, s in enumerate(slides, 1):
        if not isinstance(s, dict):
            errors.append(f"slide {idx}: must be object")
            continue
        st = str(s.get("type") or "content").lower()
        if st in ("content", "evidence_table"):
            takeaway = str(s.get("takeaway") or "").strip()
            citations = s.get("citations") or []
            if isinstance(citations, str):
                citations = [citations]
            citations = [str(c).strip() for c in citations if str(c).strip()]
            if not takeaway:
                errors.append(f"slide {idx} ({st}): missing takeaway")
            if not citations:
                errors.append(f"slide {idx} ({st}): missing citations")

        if st == "evidence_table":
            table = s.get("table") or {}
            headers = table.get("headers")
            rows = table.get("rows")
            if not isinstance(headers, list) or len(headers) < 2:
                errors.append(f"slide {idx} (evidence_table): table.headers invalid")
            if not isinstance(rows, list) or len(rows) < 1:
                errors.append(f"slide {idx} (evidence_table): table.rows invalid")

    return errors


def _safe_title_to_filename(title: str) -> str:
    base = (title or "presentation")[:60].replace("/", "-").strip()
    return base or "presentation"


def _load_template() -> Presentation:
    template_path = Path("app/templates/minimal_medical.pptx")
    if template_path.exists():
        return Presentation(str(template_path))
    return Presentation()


def _force_widescreen(prs: Presentation) -> None:
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def _normalize_outline(outline: Any) -> dict[str, Any]:
    """Accept dict or raw string; try JSON parsing when needed."""
    if isinstance(outline, dict):
        return outline
    if isinstance(outline, str):
        # Try to extract JSON if model wrapped it
        s = outline.strip()
        try:
            return json.loads(s)
        except Exception:
            # naive fallback: return as raw
            return {"title": "Presentation", "raw": s, "slides": []}
    return {"title": "Presentation", "slides": []}


async def generate_presentation(outline: dict[str, Any] | str, references: list[str], project_id: UUID) -> dict[str, Any]:
    """Genera PPTX desde outline JSON y lo guarda via StorageManager.

    Outline esperado:
    {
      "title": "...",
      "slides": [
        {"type": "title", "title": "...", "subtitle": "..."},
        {"type": "section", "title": "..."},
        {"type": "content", "title": "...", "bullets": ["..."], "notes": "..."},
        {"type": "references"}
      ]
    }
    """

    o = _normalize_outline(outline)
    prs = _load_template()
    _force_widescreen(prs)

    slides_data = o.get("slides", []) if isinstance(o, dict) else []
    if not isinstance(slides_data, list):
        slides_data = []

    # Validate to reduce generic decks and ensure traceability
    errors = validate_outline(o if isinstance(o, dict) else {})
    if errors:
        raise ValueError("Outline validation failed: " + "; ".join(errors[:6]))

    if len(slides_data) > MAX_SLIDES:
        raise ValueError(f"Máximo {MAX_SLIDES} slides permitidos")

    for slide_data in slides_data:
        if not isinstance(slide_data, dict):
            continue
        slide_type = (slide_data.get("type") or "content").lower()
        if slide_type == "title":
            _create_title_slide(prs, slide_data)
        elif slide_type in ("section", "divider"):
            _create_section_slide(prs, slide_data)
        elif slide_type == "references":
            _create_references_slide(prs, references)
        elif slide_type in ("evidence_table", "key_evidence"):
            _create_evidence_table_slide(prs, slide_data)
        else:
            _create_content_slide(prs, slide_data)

    # Ensure we always have at least 1 slide
    if len(prs.slides) == 0:
        _create_title_slide(prs, {"title": o.get("title", "Presentation"), "subtitle": ""})
        _create_references_slide(prs, references)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    filename = f"{_safe_title_to_filename(o.get('title', 'Presentation'))}_{datetime.now():%Y%m%d}.pptx"
    file_info = await storage_manager.save_presentation(bio.getvalue(), filename, project_id)

    return {
        "filename": file_info["filename"],
        "file_path": file_info["file_path"],
        "slide_count": len(prs.slides),
        "size_kb": file_info["size_kb"],
    }


def _apply_title_style(shape) -> None:
    if not shape or not getattr(shape, "text_frame", None):
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0B, 0x2A, 0x4A)


def _create_title_slide(prs: Presentation, data: dict[str, Any]) -> None:
    layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None
    slide = prs.slides.add_slide(layout) if layout else prs.slides.add_slide(prs.slide_layouts[0])

    title = slide.shapes.title
    subtitle_text = data.get("subtitle") or data.get("content") or ""

    if title:
        title.text = data.get("title", "")
        _apply_title_style(title)

    # Subtitle placeholder if exists
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    else:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.2))
        tf = box.text_frame
        tf.text = subtitle_text
        p = tf.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _create_section_slide(prs: Presentation, data: dict[str, Any]) -> None:
    # Use a blank-ish layout if possible
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Clear existing placeholders text
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            shp.text_frame.clear()

    title_text = data.get("title", "")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0B, 0x2A, 0x4A)


def _create_content_slide(prs: Presentation, data: dict[str, Any]) -> None:
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_text = data.get("title", "")
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        _apply_title_style(slide.shapes.title)

    # Body placeholder
    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    bullets = data.get("bullets") or data.get("content") or []
    if isinstance(bullets, str):
        bullets = [bullets]

    # Takeaway and citations (non-generic requirement)
    takeaway = (data.get("takeaway") or "").strip()
    citations = data.get("citations") or []
    if isinstance(citations, str):
        citations = [citations]

    if body and getattr(body, "text_frame", None):
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        # Soft warning in speaker notes
        if len(list(bullets)) > 6:
            existing_notes = slide.notes_slide.notes_text_frame.text or ""
            slide.notes_slide.notes_text_frame.text = (
                existing_notes.rstrip() + "\n\nWARNING: demasiados bullets (>6). Considera simplificar."
            ).strip()

        for i, bullet in enumerate(list(bullets)[:6]):
            if i == 0:
                tf.text = str(bullet)
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = str(bullet)
            p.level = 0
            p.font.name = "Calibri"
            p.font.size = Pt(24 if i == 0 else 22)
            p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
    else:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.3), Inches(4.8))
        tf = box.text_frame
        tf.clear()
        for i, bullet in enumerate(list(bullets)[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(bullet)
            p.font.name = "Calibri"
            p.font.size = Pt(24 if i == 0 else 22)

    # Footer: takeaway + citations
    footer_y = Inches(6.9)
    footer = slide.shapes.add_textbox(Inches(0.8), footer_y, Inches(11.7), Inches(0.5))
    ft = footer.text_frame
    ft.clear()
    p = ft.paragraphs[0]
    parts = []
    if takeaway:
        parts.append(f"Takeaway: {takeaway}")
    if citations:
        parts.append("Refs: " + "; ".join([str(c).strip() for c in list(citations)[:3] if str(c).strip()]))
    p.text = " | ".join(parts) if parts else ""
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Notes: always include citations + any provided notes
    notes = (data.get("notes") or "").strip()
    if citations:
        notes = (notes + "\n\nCITATIONS:\n" + "\n".join([f"- {c}" for c in citations])).strip()
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _create_evidence_table_slide(prs: Presentation, data: dict[str, Any]) -> None:
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = data.get("title", "Key Evidence")
        _apply_title_style(slide.shapes.title)

    table_spec = data.get("table") or {}
    headers = table_spec.get("headers") or []
    rows = table_spec.get("rows") or []

    # Coerce to strings
    headers = [str(h) for h in headers]
    rows = [[str(c) for c in r] for r in rows if isinstance(r, list)]

    cols = max(2, len(headers))
    body_rows = max(1, len(rows))

    # Table area
    x = Inches(0.8)
    y = Inches(1.6)
    w = Inches(11.7)
    h = Inches(4.9)

    shape = slide.shapes.add_table(body_rows + 1, cols, x, y, w, h)
    tbl = shape.table

    # Headers
    for j in range(cols):
        cell = tbl.cell(0, j)
        cell.text = headers[j] if j < len(headers) else ""
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0B, 0x2A, 0x4A)

    # Body
    for i in range(body_rows):
        row = rows[i] if i < len(rows) else [""] * cols
        for j in range(cols):
            cell = tbl.cell(i + 1, j)
            cell.text = row[j] if j < len(row) else ""
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Calibri"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

    # Footer like content slides
    takeaway = (data.get("takeaway") or "").strip()
    citations = data.get("citations") or []
    if isinstance(citations, str):
        citations = [citations]

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.5))
    ft = footer.text_frame
    ft.clear()
    p = ft.paragraphs[0]
    parts = []
    if takeaway:
        parts.append(f"Takeaway: {takeaway}")
    if citations:
        parts.append("Refs: " + "; ".join([str(c).strip() for c in list(citations)[:3] if str(c).strip()]))
    p.text = " | ".join(parts) if parts else ""
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    notes = (data.get("notes") or "").strip()
    if citations:
        notes = (notes + "\n\nCITATIONS:\n" + "\n".join([f"- {c}" for c in citations])).strip()
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _create_references_slide(prs: Presentation, references: list[str]) -> None:
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = "Referencias"
        _apply_title_style(slide.shapes.title)

    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if body and getattr(body, "text_frame", None):
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, ref in enumerate(references[:25], 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = ref
            p.level = 0
            p.font.name = "Calibri"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
    else:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.3), Inches(5.2))
        tf = box.text_frame
        tf.clear()
        for ref in references[:25]:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text else tf.add_paragraph()
            p.text = ref
            p.font.name = "Calibri"
            p.font.size = Pt(14)
